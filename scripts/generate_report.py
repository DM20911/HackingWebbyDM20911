#!/usr/bin/env python3
"""
generate_report.py — Generador multi-formato de informes ofensivos.

Uso:
    python3 generate_report.py --input findings.json --format docx --output report.docx
    python3 generate_report.py --input findings.json --format md   --output report.md
    python3 generate_report.py --input findings.json --format pdf  --output report.pdf
    python3 generate_report.py --input findings.json --format html --output report.html
    python3 generate_report.py --input findings.json --format txt  --output report.txt
    python3 generate_report.py --input findings.json --format pptx --output presentation.pptx
    python3 generate_report.py --input findings.json --format all  --output ./out/

Input JSON schema: ver `references/reporting.md` y ejemplo en `assets/example_findings.json`.

Texto justificado por defecto en docx/pdf/html. Sin guiones largos en cuerpo.
"""

import argparse
import json
import os
import re
import sys
from datetime import date


SEVERITY_ORDER = {"crítica": 0, "critica": 0, "critical": 0,
                  "alta": 1, "alto": 1, "high": 1,
                  "media": 2, "medio": 2, "medium": 2,
                  "baja": 3, "bajo": 3, "low": 3,
                  "info": 4, "informativa": 4, "informative": 4}

SEVERITY_COLOR = {0: "8B0000", 1: "FF0000", 2: "FF8C00", 3: "00B050", 4: "808080"}


# ─── Utilidades de texto ──────────────────────────────────────────────────────
def strip_dashes(s):
    """Quita em/en dashes en cuerpo (no afecta títulos ni código)."""
    if not isinstance(s, str):
        return s
    return re.sub(r"\s*[—–]\s*", ", ", s)


def severity_key(f):
    sev = (f.get("severidad") or f.get("severity") or "").strip().lower()
    return SEVERITY_ORDER.get(sev, 99)


def normalize(report):
    """Sanitiza cuerpo y ordena hallazgos por severidad."""
    for k in ("introduccion", "resumen_ejecutivo", "alcance", "metodologia",
              "limitaciones", "conclusion"):
        if k in report:
            report[k] = strip_dashes(report[k])
    for f in report.get("findings", []):
        for k in ("descripcion", "impacto", "remediacion", "riesgo_negocio",
                  "evidencia"):
            if k in f:
                f[k] = strip_dashes(f[k])
        # PoC steps: si es objeto, sanitize descripcion
        for step in (f.get("poc") or []):
            if isinstance(step, dict):
                if "descripcion" in step:
                    step["descripcion"] = strip_dashes(step["descripcion"])
    report["findings"] = sorted(report.get("findings", []), key=severity_key)
    return report


# ─── Markdown ─────────────────────────────────────────────────────────────────
def render_md(report):
    meta = report.get("meta", {})
    out = []
    out.append(f"# {meta.get('titulo', 'Informe de Auditoría de Seguridad')}\n")
    out.append(f"**Cliente:** {meta.get('cliente', '')}  ")
    out.append(f"**Tipo:** {meta.get('tipo', '')}  ")
    out.append(f"**Fecha:** {meta.get('fecha', date.today().isoformat())}  ")
    out.append(f"**Clasificación:** {meta.get('clasificacion', 'Confidencial')}  ")
    out.append(f"**Versión:** {meta.get('version', '1.0')}  ")
    out.append(f"**Autor:** {meta.get('autor', '')}\n")

    secs = [
        ("Introducción", report.get("introduccion")),
        ("Resumen Ejecutivo", report.get("resumen_ejecutivo")),
        ("Alcance", report.get("alcance")),
        ("Metodología", report.get("metodologia")),
    ]
    for title, body in secs:
        if body:
            out.append(f"\n## {title}\n\n{body}\n")

    # Resumen de hallazgos
    out.append("\n## Resumen de Hallazgos\n")
    out.append("| ID | Vulnerabilidad | Severidad | CVSS 3.1 | CVSS 4.0 | CWE | Estado |")
    out.append("|----|----------------|-----------|----------|----------|-----|--------|")
    for i, f in enumerate(report.get("findings", []), 1):
        out.append(f"| H{i} | {f.get('titulo','')} | {f.get('severidad','')} | "
                   f"{f.get('cvss31_score','')} | {f.get('cvss40_score','')} | "
                   f"{f.get('cwe','')} | {f.get('estado','Confirmado')} |")

    # Hallazgos detallados — cada uno como tabla 2 columnas (campo / valor)
    out.append("\n## Hallazgos Técnicos\n")
    for i, f in enumerate(report.get("findings", []), 1):
        sev = f.get('severidad','')
        sev_color = SEVERITY_COLOR.get(severity_key(f), "808080")
        # Cada hallazgo desde el #2 arranca en página nueva (HTML/PDF)
        if i > 1:
            out.append('\n<div class="pagebreak"></div>\n')
        # Cabecera con color por severidad usando HTML inline (GitHub renderea)
        out.append(f"\n### H{i} — {f.get('titulo','')}\n")
        # Tabla principal del hallazgo (todo MENOS la PoC)
        rows = []
        rows.append(("Severidad", f"<span style=\"background:#{sev_color};color:#fff;padding:2px 8px;border-radius:3px;font-weight:bold\">{sev}</span>"))
        if f.get("cvss31_vector") or f.get("cvss31_score"):
            rows.append(("CVSS 3.1", f"<strong>{f.get('cvss31_score','')}</strong>&nbsp;&nbsp;<span style=\"font-family:Menlo,Consolas,monospace;font-size:0.78em;white-space:normal;word-break:break-all\">{f.get('cvss31_vector','')}</span>"))
        if f.get("cvss40_vector") or f.get("cvss40_score"):
            rows.append(("CVSS 4.0", f"<strong>{f.get('cvss40_score','')}</strong>&nbsp;&nbsp;<span style=\"font-family:Menlo,Consolas,monospace;font-size:0.78em;white-space:normal;word-break:break-all\">{f.get('cvss40_vector','')}</span>"))
        if f.get("cwe"):     rows.append(("CWE", f['cwe']))
        if f.get("capec"):   rows.append(("CAPEC", f['capec']))
        if f.get("owasp"):   rows.append(("OWASP", f['owasp']))
        if f.get("attack"):  rows.append(("MITRE ATT&CK", f['attack']))
        if f.get("activo"):  rows.append(("Activo afectado", f"`{f['activo']}`"))
        if f.get("parametro"): rows.append(("Parámetro vulnerable", f"`{f['parametro']}`"))
        rows.append(("Descripción", (f.get('descripcion','') or '').replace("\n"," <br> ")))
        rows.append(("Impacto", (f.get('impacto','') or '').replace("\n"," <br> ")))
        if f.get("riesgo_negocio"):
            rows.append(("Riesgo de negocio", f['riesgo_negocio'].replace("\n"," <br> ")))
        rows.append(("Confidence", f.get('confidence','High')))
        rows.append(("Validación manual", f.get('validacion_manual','Sí')))
        rows.append(("Probabilidad", f.get('probabilidad','Alta')))
        if f.get("evidencia"):
            rows.append(("Evidencia técnica", f['evidencia'].replace("\n"," <br> ")))
        if f.get("remediacion"):
            rows.append(("Remediación", f['remediacion'].replace("\n"," <br> ")))
        if f.get("referencias"):
            refs = " <br> ".join(f"• {r}" for r in f["referencias"])
            rows.append(("Referencias", refs))
        rows.append(("Estado", f.get('estado','Confirmado')))
        out.append("")
        out.append("| Concepto | Descripción |")
        out.append("|----------|-------------|")
        for k, v in rows:
            out.append(f"| **{k}** | {v} |")
        # PoC fuera de la tabla (formato actual con bloques de código)
        out.append("\n#### Prueba de Concepto (PoC)\n")
        for n, step in enumerate(f.get("poc") or [], 1):
            if isinstance(step, str):
                out.append(f"\n**Paso {n}:** {step}\n")
            elif isinstance(step, dict):
                desc = step.get("descripcion") or step.get("desc") or ""
                cmd = step.get("cmd") or step.get("comando") or ""
                resp = step.get("output") or step.get("respuesta") or ""
                out.append(f"\n**Paso {n}:** {desc}\n")
                if cmd:
                    out.append(f"\n```bash\n{cmd}\n```\n")
                if resp:
                    out.append(f"\n```\n{resp}\n```\n")

    # Attack chains
    if report.get("attack_chains"):
        out.append("\n## Attack Chains\n")
        for c in report["attack_chains"]:
            out.append(f"\n### {c.get('nombre','Chain')}\n")
            out.append(f"**Componentes:** {', '.join(c.get('componentes',[]))}\n")
            out.append(f"**CVSS agregado:** {c.get('cvss_agregado','')}\n\n{c.get('descripcion','')}\n")

    # Matriz de riesgo
    counts = {}
    for f in report.get("findings", []):
        sev = f.get("severidad", "Info")
        counts[sev] = counts.get(sev, 0) + 1
    if counts:
        out.append("\n## Matriz de Riesgo\n")
        # ── Matriz 5x5 estilo heat-map (Probabilidad × Impacto) ────────────
        # Y (Probabilidad, top→bottom de alta a baja):
        Y_LABELS = ["Casi seguro", "Probable", "Posible", "Improbable", "Raro"]
        # X (Impacto, left→right de bajo a alto):
        X_LABELS = ["Insignificante", "Menor", "Moderado", "Mayor", "Catastrófico"]
        # Score = (5 - y) * (x + 1) con y,x ∈ 0..4 → rango 1..25
        # Verde 1-5, amarillo 6-10, naranja 11-15, rojo 16-25
        def cell_color(y, x):
            score = (5 - y) * (x + 1)
            if score <= 5:  return "00B050"   # verde
            if score <= 10: return "FFD600"   # amarillo
            if score <= 15: return "FF8C00"   # naranja
            return "C00000"                   # rojo intenso

        # Mapear cada finding a (y, x) usando probabilidad + severidad
        def prob_idx(s):
            s = (s or "").strip().lower()
            if s in ("casi seguro","muy alta","muy probable"): return 0
            if s in ("probable","alta","alto","high"):          return 1
            if s in ("posible","media","medio","medium"):       return 2
            if s in ("improbable","baja","bajo","low"):         return 3
            if s in ("raro","muy baja"):                        return 4
            return 2  # default Posible
        def imp_idx(sev):
            sev = (sev or "").strip().lower()
            if sev in ("crítica","critica","critical"): return 4
            if sev in ("alta","alto","high"):           return 3
            if sev in ("media","medio","medium"):       return 2
            if sev in ("baja","bajo","low"):            return 1
            return 0

        grid = [[0]*5 for _ in range(5)]
        for f in report.get("findings", []):
            y = prob_idx(f.get("probabilidad"))
            x = imp_idx(f.get("severidad"))
            grid[y][x] += 1

        # Render del heat-map como tabla 6x6 (col 0 = labels Y, row 0 = labels X)
        # Layout flex: etiqueta Y a la izquierda (fuera del área de la tabla)
        # + tabla + etiqueta X debajo.
        out.append('<div class="riskmatrix-outer">')
        out.append('<div class="rm-y-axis">Probabilidad ▲</div>')
        out.append('<div class="riskmatrix-wrap">')
        out.append('<table class="riskmatrix"><tbody>')
        # Encabezado: esquina vacía + impactos
        out.append('<tr><th class="rm-corner"></th>')
        for xl in X_LABELS:
            out.append(f'<th class="rm-xhead">{xl}</th>')
        out.append('</tr>')
        # Filas
        for y, yl in enumerate(Y_LABELS):
            out.append(f'<tr><th class="rm-yhead">{yl}</th>')
            for x in range(5):
                color = cell_color(y, x)
                n = grid[y][x]
                badge = f'<span class="rm-count">{n}</span>' if n > 0 else ''
                out.append(f'<td class="rm-cell" style="background:#{color}">{badge}</td>')
            out.append('</tr>')
        out.append('</tbody></table>')
        out.append('<div class="rm-x-axis">Impacto ▶</div>')
        out.append('</div>')   # cierra riskmatrix-wrap
        out.append('</div>')   # cierra riskmatrix-outer

        # Leyenda de colores
        out.append('<div class="rm-legend">'
                   '<span class="rm-leg" style="background:#00B050">Bajo</span>'
                   '<span class="rm-leg" style="background:#FFD600;color:#000">Medio</span>'
                   '<span class="rm-leg" style="background:#FF8C00">Alto</span>'
                   '<span class="rm-leg" style="background:#C00000">Crítico</span>'
                   '</div>')

    if report.get("timeline"):
        out.append("\n## Timeline de Testing\n")
        for t in report["timeline"]:
            out.append(f"- {t}")

    if report.get("limitaciones"):
        out.append(f"\n## Limitaciones\n\n{report['limitaciones']}\n")

    if report.get("conclusion"):
        # Conclusión siempre en página nueva
        out.append('\n<div class="pagebreak"></div>\n')
        out.append(f"\n## Conclusión General\n\n{report['conclusion']}\n")

    if report.get("anexos"):
        out.append("\n## Anexos\n")
        for k, v in report["anexos"].items():
            out.append(f"\n### {k}\n\n{v}\n")

    return "\n".join(out)


# ─── TXT ──────────────────────────────────────────────────────────────────────
def render_txt(report):
    md = render_md(report)
    # Sacar marcas markdown obvias
    txt = re.sub(r"^#{1,6}\s+", "", md, flags=re.M)
    txt = re.sub(r"\*\*(.+?)\*\*", r"\1", txt)
    txt = re.sub(r"`(.+?)`", r"\1", txt)
    txt = re.sub(r"^[\-\*]\s+", "  • ", txt, flags=re.M)
    return txt


# ─── HTML ─────────────────────────────────────────────────────────────────────
def render_html(report):
    try:
        import markdown
    except ImportError:
        os.system(f"{sys.executable} -m pip install --quiet markdown")
        import markdown
    md_body = render_md(report)
    html_body = markdown.markdown(md_body, extensions=["tables", "fenced_code"])
    style = """
    <style>
      body { font-family: -apple-system, Segoe UI, Roboto, Helvetica, Arial, sans-serif;
             max-width: 900px; margin: 2em auto; padding: 0 2em;
             text-align: justify; line-height: 1.55; color: #1a1a1a; }
      h1, h2, h3 { color: #1F2937; text-align: left; }
      h1 { border-bottom: 3px solid #8B0000; padding-bottom: .3em; }
      h2 { border-bottom: 1px solid #ccc; padding-bottom: .2em; margin-top: 2em; }
      h3 { margin-top: 1.6em; }
      table { border-collapse: collapse; width: 100%; margin: 1em 0; table-layout: fixed; }
      th, td { border: 1px solid #ccc; padding: .5em .8em; text-align: justify; vertical-align: top; word-break: break-word; overflow-wrap: anywhere; hyphens: auto; }
      th { background: #F2F2F2; text-align: left; font-weight: bold; }
      /* Tablas de hallazgos: columna izquierda angosta, derecha justificada */
      h3 + table th:first-child, h3 + table td:first-child { width: 22%; font-weight: bold; background: #F7F7F7; text-align: left; }
      h3 + table th:last-child,  h3 + table td:last-child  { width: 78%; }
      /* Tabla "Resumen de Hallazgos" (la primera tras un h2): compacta,
         tipografía pequeña, columnas que se ajustan al contenido,
         encabezados sin salto de línea, valores cortos centrados. */
      h2 + table { font-size: 0.78em; table-layout: auto; width: 100%; margin-top: .4em; }
      h2 + table th { white-space: nowrap; padding: .35em .55em; text-align: left; background: #F2F2F2; }
      h2 + table td { padding: .3em .55em; vertical-align: top; text-align: left; }
      /* Columnas cortas: ID, Severidad, CVSS 3.1, CVSS 4.0, Estado — sin wrap, centradas */
      h2 + table td:nth-child(1),
      h2 + table td:nth-child(3),
      h2 + table td:nth-child(4),
      h2 + table td:nth-child(5),
      h2 + table td:nth-child(7) { white-space: nowrap; text-align: center; }
      /* Matriz de Riesgo 5x5 — heat-map clásico Probabilidad × Impacto */
      .riskmatrix-outer { display: flex; align-items: center; gap: 10px;
                          max-width: 640px; margin: 1.5em auto;
                          page-break-inside: avoid; break-inside: avoid; }
      .rm-y-axis { writing-mode: vertical-rl; transform: rotate(180deg);
                   font-weight: 700; color: #1F2937; letter-spacing: 1px;
                   font-size: .8em; flex-shrink: 0; white-space: nowrap;
                   padding: 0 2px; }
      .riskmatrix-wrap { flex: 1; text-align: center; }
      .rm-x-axis { margin-top: .6em; font-weight: 700; color: #1F2937;
                   letter-spacing: .5px; font-size: .8em; }
      table.riskmatrix { border-collapse: separate; border-spacing: 3px;
                         margin: 0 auto; table-layout: fixed; width: 100%;
                         page-break-inside: avoid; break-inside: avoid; }
      .rm-corner { background: transparent !important; border: none !important; }
      .rm-xhead { background: #1F2937 !important; color: #fff !important;
                  padding: .35em .25em !important; font-size: .65em !important;
                  text-align: center !important; font-weight: 700;
                  border-radius: 4px; white-space: normal; line-height: 1.15; }
      .rm-yhead { background: #1F2937 !important; color: #fff !important;
                  padding: .35em .45em !important; font-size: .65em !important;
                  text-align: right !important; font-weight: 700;
                  border-radius: 4px; white-space: nowrap; }
      .rm-cell  { height: 44px; border-radius: 4px; position: relative;
                  vertical-align: middle; text-align: center !important;
                  border: none !important; }
      .rm-count { display: inline-block; min-width: 22px; height: 22px; line-height: 22px;
                  background: rgba(255,255,255,.95); color: #1F2937;
                  font-weight: 800; font-size: .95em; border-radius: 50%;
                  box-shadow: 0 1px 3px rgba(0,0,0,.25); padding: 0 5px; }
      .rm-legend { display: flex; gap: 8px; justify-content: center;
                   margin: .6em auto 1.4em; font-size: .78em; }
      .rm-leg { color: #fff; padding: 3px 14px; border-radius: 12px;
                font-weight: 700; letter-spacing: .3px; }
      /* Salto de página por hallazgo (HTML/PDF). Soporte legacy + moderno */
      .pagebreak { page-break-before: always; break-before: page; }
      h3 { page-break-after: avoid; break-after: avoid; }
      table { page-break-inside: auto; break-inside: auto; }
      tr { page-break-inside: avoid; break-inside: avoid; }
      pre { page-break-inside: avoid; break-inside: avoid; }
      pre { background: #1E1E1E; color: #33FF66; padding: 1em; border-radius: 4px;
            text-align: left; font-family: Menlo, Consolas, monospace;
            white-space: pre-wrap; word-break: break-all; overflow-wrap: anywhere; }
      code { font-family: Menlo, Consolas, monospace; background: #F5F5F5;
             padding: .1em .3em; border-radius: 3px; }
      pre code { background: transparent; color: inherit; padding: 0; }
      blockquote { border-left: 4px solid #8B0000; margin: 1em 0; padding-left: 1em;
                   color: #555; }
    </style>
    """
    title = report.get("meta", {}).get("titulo", "Informe de Auditoría")
    return f"<!doctype html><html><head><meta charset='utf-8'><title>{title}</title>{style}</head><body>{html_body}</body></html>"


# ─── DOCX ─────────────────────────────────────────────────────────────────────
def render_docx(report, path):
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        os.system(f"{sys.executable} -m pip install --quiet python-docx")
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

    doc = Document()

    # Márgenes 2cm
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2)
        section.right_margin = Cm(2)

    def set_shading(cell, fill):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), fill)
        tcPr.append(shd)

    # Numerador global para captions de PoC ("Imagen Nº 1", "Imagen Nº 2", ...)
    fig_counter = {"n": 0}

    def add_caption(text):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.italic = True
        run.font.name = "Arial"
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor.from_string("555555")
        p.paragraph_format.space_after = Pt(6)

    def add_para(text, bold=False, size=10.5, justify=True, color=None, space_after=6):
        p = doc.add_paragraph()
        if justify:
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        run = p.add_run(text)
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.bold = bold
        if color:
            run.font.color.rgb = RGBColor.from_string(color)
        p.paragraph_format.space_after = Pt(space_after)
        return p

    def add_heading(text, level=1):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.name = "Arial"
            run.font.color.rgb = RGBColor.from_string("1F2937")
        # Mantener el título junto al siguiente párrafo y evitar viudas:
        # si queda muy abajo de la página, Word lo empuja a la siguiente.
        h.paragraph_format.keep_with_next = True
        h.paragraph_format.keep_together = True
        return h

    def add_code_block(text, kind="cmd"):
        # Bloque oscuro estilo terminal. Wrap forzado para que líneas largas
        # (URLs, payloads) no se salgan del margen del documento.
        for i, line in enumerate(text.splitlines() or [""]):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            prefix = "$ " if (kind == "cmd" and i == 0) else ("  " if kind == "cmd" else "")
            run = p.add_run(prefix + (line or " "))
            run.font.name = "Menlo"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor.from_string("33FF66" if kind == "cmd" else "BFFFCB")
            p.paragraph_format.space_after = Pt(0)
            # Indent leve para visual de bloque y limitar ancho efectivo
            p.paragraph_format.left_indent = Cm(0.2)
            p.paragraph_format.right_indent = Cm(0.2)
            pPr = p._p.get_or_add_pPr()
            # Shading oscuro
            shd = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear'); shd.set(qn('w:color'), 'auto'); shd.set(qn('w:fill'), '1E1E1E')
            pPr.append(shd)
            # Forzar wrap dentro de palabras Latinas largas (URLs, base64, payloads)
            ww = OxmlElement('w:wordWrap')
            ww.set(qn('w:val'), '0')   # 0 = puede partir palabras para evitar overflow
            pPr.append(ww)

    meta = report.get("meta", {})
    add_heading(meta.get("titulo", "Informe de Auditoría de Seguridad"), level=0)
    add_para(f"Cliente: {meta.get('cliente','')}", bold=True, size=11, justify=False, space_after=2)
    add_para(f"Tipo: {meta.get('tipo','')}", justify=False, space_after=2)
    add_para(f"Fecha: {meta.get('fecha', date.today().isoformat())}", justify=False, space_after=2)
    add_para(f"Clasificación: {meta.get('clasificacion','Confidencial')}", justify=False, space_after=2)
    add_para(f"Versión: {meta.get('version','1.0')}", justify=False, space_after=2)
    add_para(f"Autor: {meta.get('autor','')}", justify=False, space_after=12)

    for title, body in [("Introducción", report.get("introduccion")),
                        ("Resumen Ejecutivo", report.get("resumen_ejecutivo")),
                        ("Alcance", report.get("alcance")),
                        ("Metodología", report.get("metodologia"))]:
        if body:
            add_heading(title, 1)
            add_para(body)

    # Tabla resumen
    add_heading("Resumen de Hallazgos", 1)
    findings = report.get("findings", [])
    if findings:
        table = doc.add_table(rows=1, cols=7)
        table.style = "Light Grid"
        hdr = ["ID", "Vulnerabilidad", "Severidad", "CVSS 3.1", "CVSS 4.0", "CWE", "Estado"]
        for i, h in enumerate(hdr):
            c = table.rows[0].cells[i]
            c.text = h
            for r in c.paragraphs[0].runs:
                r.bold = True
            set_shading(c, "F2F2F2")
        for i, f in enumerate(findings, 1):
            row = table.add_row().cells
            row[0].text = f"H{i}"
            row[1].text = f.get("titulo", "")
            row[2].text = f.get("severidad", "")
            row[3].text = str(f.get("cvss31_score", ""))
            row[4].text = str(f.get("cvss40_score", ""))
            row[5].text = f.get("cwe", "")
            row[6].text = f.get("estado", "Confirmado")
            sev_color = SEVERITY_COLOR.get(severity_key(f), "808080")
            set_shading(row[2], sev_color)

    # ─── Hallazgos como TABLAS (estilo moderno) ─────────────────────────────
    # Cada hallazgo: tabla 2 columnas (Campo / Valor) con encabezado oscuro,
    # fila de severidad con color por criticidad, etiquetas grises.
    # La PoC va FUERA de la tabla (manteniendo bloques de código tipo terminal).
    def cell_text(cell, text, bold=False, color=None, size=10.5, font="Arial", align=None):
        # Renderiza el contenido de la celda en UN solo párrafo, dejando que Word
        # haga el wrap natural según el ancho de la celda (sin "br hacia abajo"
        # forzados por cada \n del input). Los \n del input se convierten en
        # line-breaks suaves dentro del mismo párrafo para conservar la lectura.
        cell.text = ""
        p = cell.paragraphs[0]
        if align is not None: p.alignment = align
        s = str(text or "")
        parts = s.split("\n")
        for idx, part in enumerate(parts):
            if idx > 0:
                # soft break dentro del mismo párrafo
                br = OxmlElement('w:br'); p._p.find(qn('w:r')) if False else None
                run = p.add_run()
                run._r.append(OxmlElement('w:br'))
            run = p.add_run(part)
            run.font.name = font
            run.font.size = Pt(size)
            run.bold = bold
            if color:
                run.font.color.rgb = RGBColor.from_string(color)

    add_heading("Hallazgos Técnicos", 1)
    LABEL_W = 2400  # DXA
    VALUE_W = 6800
    for i, f in enumerate(findings, 1):
        if i > 1:
            doc.add_page_break()
        sev = f.get('severidad','')
        sev_color = SEVERITY_COLOR.get(severity_key(f), "808080")

        # Construir filas (label, value, kind) — todo MENOS la PoC
        # kind: "text" = Arial 10.5 justificado
        #       "mono" = Menlo 9 justificado (host/parámetro)
        #       "cvss" = score Arial bold + vector Menlo pequeño en MISMA línea
        rows = []
        if f.get("cvss31_vector") or f.get("cvss31_score"):
            rows.append(("CVSS 3.1", (f.get('cvss31_score',''), f.get('cvss31_vector','')), "cvss"))
        if f.get("cvss40_vector") or f.get("cvss40_score"):
            rows.append(("CVSS 4.0", (f.get('cvss40_score',''), f.get('cvss40_vector','')), "cvss"))
        for key, label in [("cwe","CWE"),("capec","CAPEC"),("owasp","OWASP"),
                            ("attack","MITRE ATT&CK"),
                            ("activo","Activo afectado"),
                            ("parametro","Parámetro vulnerable")]:
            if f.get(key): rows.append((label, f[key], "mono" if key in ("activo","parametro") else "text"))
        rows.append(("Descripción", f.get('descripcion',''), "text"))
        rows.append(("Impacto", f.get('impacto',''), "text"))
        if f.get("riesgo_negocio"): rows.append(("Riesgo de negocio", f['riesgo_negocio'], "text"))
        rows.append(("Confidence", f.get('confidence','High'), "text"))
        rows.append(("Validación manual", f.get('validacion_manual','Sí'), "text"))
        rows.append(("Probabilidad", f.get('probabilidad','Alta'), "text"))
        if f.get("evidencia"): rows.append(("Evidencia técnica", f['evidencia'], "text"))
        if f.get("remediacion"): rows.append(("Remediación", f['remediacion'], "text"))
        if f.get("referencias"):
            rows.append(("Referencias", "\n".join(f"• {r}" for r in f["referencias"]), "text"))
        rows.append(("Estado", f.get('estado','Confirmado'), "text"))

        # Crear tabla con cabecera oscura + fila de severidad coloreada
        table = doc.add_table(rows=2 + len(rows), cols=2)
        table.style = "Light Grid"
        table.autofit = False
        # Columna izquierda angosta para dar más aire al texto de la derecha.
        # Total útil A4 con márgenes 2 cm ~= 17 cm.
        LEFT_W = Cm(3.2)
        RIGHT_W = Cm(13.8)
        table.columns[0].width = LEFT_W
        table.columns[1].width = RIGHT_W
        for row in table.rows:
            row.cells[0].width = LEFT_W
            row.cells[1].width = RIGHT_W
            # Evitar que una fila se parta entre dos páginas si cabe entera
            trPr = row._tr.get_or_add_trPr()
            cant_split = OxmlElement('w:cantSplit'); trPr.append(cant_split)

        # Fila 0 — header oscuro con título del hallazgo (col-span simulado)
        hdr_left, hdr_right = table.rows[0].cells
        cell_text(hdr_left, f"Hallazgo H{i}", bold=True, color="FFFFFF", size=11)
        cell_text(hdr_right, f.get('titulo',''), bold=True, color="FFFFFF", size=11)
        set_shading(hdr_left, "1F2937")
        set_shading(hdr_right, "1F2937")

        # Fila 1 — Severidad coloreada
        sev_left, sev_right = table.rows[1].cells
        cell_text(sev_left, "Severidad", bold=True, color="FFFFFF", size=10)
        cell_text(sev_right, sev, bold=True, color="FFFFFF", size=10)
        set_shading(sev_left, sev_color)
        set_shading(sev_right, sev_color)

        # Filas restantes — etiquetas left, valores SIEMPRE justificados
        for idx, (label, value, kind) in enumerate(rows, start=2):
            r = table.rows[idx].cells
            cell_text(r[0], label, bold=True, size=10.5, align=WD_ALIGN_PARAGRAPH.LEFT)
            set_shading(r[0], "F2F2F2")
            if kind == "cvss":
                # Score Arial bold + vector Menlo pequeño en MISMA línea, justificado.
                # El vector lleva wordWrap=0 para partir en caracteres si excede el ancho.
                score, vector = value
                r[1].text = ""
                p = r[1].paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                # wrap dentro de palabras Latinas (vector CVSS es una palabra larga)
                pPr = p._p.get_or_add_pPr()
                ww = OxmlElement('w:wordWrap'); ww.set(qn('w:val'), '0'); pPr.append(ww)
                if score:
                    r0 = p.add_run(f"{score}  ")
                    r0.bold = True; r0.font.name = "Arial"; r0.font.size = Pt(10.5)
                if vector:
                    r1 = p.add_run(vector)
                    r1.font.name = "Menlo"; r1.font.size = Pt(8)
            elif kind == "mono":
                cell_text(r[1], value, font="Menlo", size=9, align=WD_ALIGN_PARAGRAPH.JUSTIFY)
            else:
                cell_text(r[1], value, size=10.5, align=WD_ALIGN_PARAGRAPH.JUSTIFY)
            set_shading(r[1], "FFFFFF")

        # Espacio después de la tabla
        doc.add_paragraph()

        # PoC FUERA de la tabla — formato actual con bloques tipo terminal
        add_heading("Prueba de Concepto (PoC)", 3)
        for n, step in enumerate(f.get("poc") or [], 1):
            if isinstance(step, str):
                add_para(f"Paso {n}: {step}", justify=True)
            elif isinstance(step, dict):
                desc = step.get("descripcion") or step.get("desc") or "agregar explicación manual"
                cmd = step.get("cmd") or step.get("comando") or ""
                resp = step.get("output") or step.get("respuesta") or ""
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                p.paragraph_format.keep_with_next = True
                p.paragraph_format.keep_together = True
                r1 = p.add_run(f"Paso {n}: "); r1.bold = True; r1.font.name = "Arial"; r1.font.size = Pt(10.5)
                r2 = p.add_run(desc); r2.font.name = "Arial"; r2.font.size = Pt(10.5)
                if cmd:
                    add_code_block(cmd, "cmd")
                    fig_counter["n"] += 1
                    add_caption(f"Imagen Nº {fig_counter['n']}: comando ejecutado en el paso {n}.")
                if resp:
                    add_code_block(resp, "output")
                    fig_counter["n"] += 1
                    add_caption(f"Imagen Nº {fig_counter['n']}: respuesta obtenida del paso {n}.")

    # Attack chains
    if report.get("attack_chains"):
        doc.add_page_break()
        add_heading("Attack Chains", 1)
        for c in report["attack_chains"]:
            add_heading(c.get("nombre", "Chain"), 2)
            add_para(f"Componentes: {', '.join(c.get('componentes',[]))}", bold=True, justify=False)
            add_para(f"CVSS agregado: {c.get('cvss_agregado','')}", justify=False)
            add_para(c.get("descripcion", ""))

    # Matriz de riesgo
    counts = {}
    for f in findings:
        sev = f.get("severidad", "Info")
        counts[sev] = counts.get(sev, 0) + 1
    if counts:
        add_heading("Matriz de Riesgo", 1)
        t = doc.add_table(rows=1, cols=2); t.style = "Light Grid"
        t.rows[0].cells[0].text = "Criticidad"; t.rows[0].cells[1].text = "Cantidad"
        for sev, n in counts.items():
            row = t.add_row().cells
            row[0].text = sev; row[1].text = str(n)

    if report.get("limitaciones"):
        add_heading("Limitaciones", 1); add_para(report["limitaciones"])
    if report.get("conclusion"):
        doc.add_page_break()
        add_heading("Conclusión General", 1); add_para(report["conclusion"])
    if report.get("anexos"):
        add_heading("Anexos", 1)
        for k, v in report["anexos"].items():
            add_heading(k, 2); add_para(v)

    doc.save(path)


# ─── PDF ──────────────────────────────────────────────────────────────────────
# Pipeline:
#   1. Intento weasyprint (mejor calidad, requiere pango/cairo nativos).
#   2. Fallback a Chrome/Chromium/Brave headless (--print-to-pdf) — siempre
#      funciona en macOS sin pelear con DYLD/SIP ni instalar libs nativas.
#   3. Fallback a wkhtmltopdf si está instalado.
#   4. Último recurso: dejar el HTML generado al lado y avisar.
def _try_weasyprint(html, path):
    # Silenciar stderr durante el probe: weasyprint imprime un troubleshooting URL
    # al fallar la carga de pango/glib en macOS antes de levantar la excepción.
    import contextlib, io
    stderr_buf = io.StringIO()
    try:
        with contextlib.redirect_stderr(stderr_buf):
            from weasyprint import HTML  # noqa
    except Exception:
        try:
            with contextlib.redirect_stderr(stderr_buf):
                os.system(f"{sys.executable} -m pip install --quiet --user weasyprint 2>/dev/null")
                from weasyprint import HTML  # noqa
        except Exception as e:
            raise RuntimeError(f"weasyprint no disponible: {e}")
    from weasyprint import HTML
    with contextlib.redirect_stderr(stderr_buf):
        HTML(string=html).write_pdf(path)


def _find_headless_browser():
    candidates = [
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Google Chrome Canary.app/Contents/MacOS/Google Chrome Canary",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Brave Browser.app/Contents/MacOS/Brave Browser",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "/usr/bin/google-chrome",
        "/usr/bin/chromium",
        "/usr/bin/chromium-browser",
    ]
    for c in candidates:
        if os.path.exists(c):
            return c
    return None


def _try_chrome_headless(html, path):
    browser = _find_headless_browser()
    if not browser:
        raise RuntimeError("no se encontró Chrome/Chromium/Brave/Edge para fallback PDF")
    import tempfile, shutil, subprocess
    with tempfile.NamedTemporaryFile("w", suffix=".html", delete=False, encoding="utf-8") as fh:
        fh.write(html)
        html_path = fh.name
    abs_pdf = os.path.abspath(path)
    try:
        subprocess.run(
            [browser, "--headless=new", "--disable-gpu", "--no-pdf-header-footer",
             f"--print-to-pdf={abs_pdf}", f"file://{html_path}"],
            check=True, capture_output=True, timeout=120,
        )
    except (subprocess.CalledProcessError, FileNotFoundError):
        subprocess.run(
            [browser, "--headless", "--disable-gpu", "--no-pdf-header-footer",
             f"--print-to-pdf={abs_pdf}", f"file://{html_path}"],
            check=True, capture_output=True, timeout=120,
        )
    finally:
        try: os.unlink(html_path)
        except Exception: pass


def _try_wkhtmltopdf(html, path):
    import shutil, tempfile, subprocess
    wk = shutil.which("wkhtmltopdf")
    if not wk:
        raise RuntimeError("wkhtmltopdf no instalado")
    with tempfile.NamedTemporaryFile("w", suffix=".html", delete=False, encoding="utf-8") as fh:
        fh.write(html)
        html_path = fh.name
    try:
        subprocess.run([wk, "--enable-local-file-access", html_path, path],
                       check=True, capture_output=True, timeout=120)
    finally:
        try: os.unlink(html_path)
        except Exception: pass


def render_pdf(report, path):
    html = render_html(report)
    errors = []
    # En macOS preferir Chrome headless: weasyprint requiere pango/cairo nativos
    # y el python protegido por SIP no respeta DYLD_LIBRARY_PATH. En Linux
    # weasyprint suele ir primero porque la calidad tipográfica es mejor.
    import platform
    if platform.system() == "Darwin":
        backends = (
            ("chrome-headless", _try_chrome_headless),
            ("weasyprint", _try_weasyprint),
            ("wkhtmltopdf", _try_wkhtmltopdf),
        )
    else:
        backends = (
            ("weasyprint", _try_weasyprint),
            ("chrome-headless", _try_chrome_headless),
            ("wkhtmltopdf", _try_wkhtmltopdf),
        )
    for backend_name, fn in backends:
        try:
            fn(html, path)
            print(f"PDF generado vía {backend_name}: {path}", file=sys.stderr)
            return
        except Exception as e:
            errors.append(f"  - {backend_name}: {e}")
            continue
    # Todos fallaron: dejar el HTML al lado y lanzar error explicativo
    html_fallback = os.path.splitext(path)[0] + ".html"
    try:
        with open(html_fallback, "w", encoding="utf-8") as fh:
            fh.write(html)
    except Exception:
        pass
    raise RuntimeError(
        "No se pudo generar el PDF con ningún backend disponible.\n"
        + "\n".join(errors)
        + f"\nHTML guardado en {html_fallback} — ábrelo en navegador y exporta a PDF manualmente."
    )


# ─── PPTX (presentación simple de resultados) ──────────────────────────────────
def render_pptx(report, path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        os.system(f"{sys.executable} -m pip install --quiet python-pptx")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    meta = report.get("meta", {})

    # Slide 1: portada
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = meta.get("titulo", "Informe de Auditoría de Seguridad")
    sub = s.placeholders[1]
    sub.text = (f"{meta.get('cliente','')}  |  {meta.get('tipo','')}\n"
                f"{meta.get('fecha', date.today().isoformat())}  |  {meta.get('clasificacion','Confidencial')}\n"
                f"Autor: {meta.get('autor','')}")

    # Slide 2: resumen ejecutivo
    if report.get("resumen_ejecutivo"):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Resumen Ejecutivo"
        s.placeholders[1].text = report["resumen_ejecutivo"][:1500]

    # Slide 3: matriz por severidad
    counts = {}
    for f in report.get("findings", []):
        sev = f.get("severidad", "Info")
        counts[sev] = counts.get(sev, 0) + 1
    if counts:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "Hallazgos por Severidad"
        rows = len(counts) + 1
        tbl = s.shapes.add_table(rows, 2, Inches(2), Inches(2), Inches(6), Inches(0.5*rows)).table
        tbl.cell(0,0).text = "Severidad"; tbl.cell(0,1).text = "Cantidad"
        for i, (sev, n) in enumerate(counts.items(), 1):
            tbl.cell(i,0).text = sev; tbl.cell(i,1).text = str(n)

    # Slide por hallazgo (top N criticos/altos)
    top = [f for f in report.get("findings", []) if severity_key(f) <= 1][:8]
    for i, f in enumerate(top, 1):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"H{i} — {f.get('titulo','')}"
        body = (f"Severidad: {f.get('severidad','')}    CVSS 3.1: {f.get('cvss31_score','')}\n"
                f"Activo: {f.get('activo','')}\n\n"
                f"Impacto: {f.get('impacto','')[:400]}\n\n"
                f"Remediación: {f.get('remediacion','')[:300]}")
        s.placeholders[1].text = body

    # Slide final: conclusión
    if report.get("conclusion"):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Conclusión General"
        s.placeholders[1].text = report["conclusion"][:1500]

    prs.save(path)


# ─── SARIF / JSON normalizado ─────────────────────────────────────────────────
def render_json(report, path):
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(report, fh, ensure_ascii=False, indent=2)


# ─── Main ─────────────────────────────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", "-i", required=True, help="JSON con findings y meta")
    ap.add_argument("--format", "-f", default="md",
                    choices=["md", "txt", "html", "docx", "pdf", "pptx", "json", "all"])
    ap.add_argument("--output", "-o", required=True)
    args = ap.parse_args()

    with open(args.input, encoding="utf-8") as fh:
        report = normalize(json.load(fh))

    if args.format == "all":
        os.makedirs(args.output, exist_ok=True)
        base = os.path.join(args.output, "report")
        with open(base + ".md", "w", encoding="utf-8") as fh: fh.write(render_md(report))
        with open(base + ".txt", "w", encoding="utf-8") as fh: fh.write(render_txt(report))
        with open(base + ".html", "w", encoding="utf-8") as fh: fh.write(render_html(report))
        render_docx(report, base + ".docx")
        try:
            render_pdf(report, base + ".pdf")
        except Exception as e:
            print(f"PDF skip: {e}", file=sys.stderr)
        render_pptx(report, base + ".pptx")
        render_json(report, base + ".json")
        print(f"OK: {args.output}/")
        return

    if args.format == "md":
        with open(args.output, "w", encoding="utf-8") as fh: fh.write(render_md(report))
    elif args.format == "txt":
        with open(args.output, "w", encoding="utf-8") as fh: fh.write(render_txt(report))
    elif args.format == "html":
        with open(args.output, "w", encoding="utf-8") as fh: fh.write(render_html(report))
    elif args.format == "docx":
        render_docx(report, args.output)
    elif args.format == "pdf":
        render_pdf(report, args.output)
    elif args.format == "pptx":
        render_pptx(report, args.output)
    elif args.format == "json":
        render_json(report, args.output)
    print(f"OK:{args.output}")


if __name__ == "__main__":
    main()
